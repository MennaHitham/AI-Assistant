import os
import logging
from typing import Optional, List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy

logger = logging.getLogger(__name__)

# Try to import MSO_SHAPE
try:
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    from pptx.util import MSO_SHAPE


class PresentationMaker:
    """
    Handles the creation of professional PowerPoint presentations.
    Supports multiple slide types: cover, content, two_column, closing.
    """

    def __init__(self, output_dir="presentations"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

        # ── Brand Colors ──────────────────────────────────────────────
        self.COLOR_BG         = RGBColor(0x1E, 0x27, 0x61)   # Deep Navy
        self.COLOR_BG_ALT     = RGBColor(0x14, 0x1A, 0x45)   # Darker Navy (cover/closing)
        self.COLOR_TITLE      = RGBColor(0xFF, 0xFF, 0xFF)   # White
        self.COLOR_BODY       = RGBColor(0xCA, 0xDC, 0xFC)   # Ice Blue
        self.COLOR_ACCENT     = RGBColor(0x00, 0xB4, 0xD8)   # Cyan Accent
        self.COLOR_ACCENT2    = RGBColor(0xFF, 0xC3, 0x00)   # Gold (closing)
        self.COLOR_MUTED      = RGBColor(0x94, 0xA3, 0xB8)   # Slate (slide numbers, captions)
        self.COLOR_CARD_BG    = RGBColor(0x25, 0x32, 0x75)   # Card background (slightly lighter)

        # ── Fonts ─────────────────────────────────────────────────────
        self.FONT_TITLE  = "Calibri"
        self.FONT_BODY   = "Calibri Light"

    # ──────────────────────────────────────────────────────────────────
    # Public API
    # ──────────────────────────────────────────────────────────────────

    def create_presentation(
        self,
        slides_data: List[Dict],
        image_paths: List[str] = None,
        filename: str = "presentation.pptx",
    ) -> Optional[str]:
        """
        Creates a professionally styled .pptx file.

        Each item in slides_data should have:
            - type:    "cover" | "content" | "two_column" | "closing"  (default: "content")
            - title:   str
            - content: list[str]  – bullet points
            - notes:   str        – speaker notes (optional)

        image_paths: optional list of local image file paths provided by the user.
            - If None or empty  → no images are used anywhere, period.
            - If provided       → images are distributed across non-cover/non-closing slides
                                  in order. Slides without a matching image render as
                                  plain content slides (no image inserted).
        """
        try:
            prs = Presentation()
            # Widescreen 16:9
            prs.slide_width  = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Validate provided image paths — keep only real, existing files
            valid_images: List[str] = []
            for p in (image_paths or []):
                if p and os.path.exists(p):
                    valid_images.append(p)
                else:
                    logger.warning(f"Image path not found, skipping: {p}")

            has_images = len(valid_images) > 0

            # Build an image assignment map:
            # Assign images only to content/two_column slides, in order.
            # cover and closing slides never receive images.
            image_iter = iter(valid_images)
            img_map: Dict[int, str] = {}   # slide index → image path
            if has_images:
                for idx, slide_info in enumerate(slides_data):
                    stype = slide_info.get("type", "content").lower()
                    if stype not in ("cover", "closing"):
                        img = next(image_iter, None)
                        if img:
                            img_map[idx] = img

            for idx, slide_info in enumerate(slides_data):
                slide_type = slide_info.get("type", "content").lower()
                img_path   = img_map.get(idx)   # None if no image assigned

                if slide_type == "cover":
                    self._add_cover_slide(prs, slide_info, idx + 1, len(slides_data))
                elif slide_type == "closing":
                    self._add_closing_slide(prs, slide_info, idx + 1, len(slides_data))
                elif img_path:
                    # User-supplied images take priority for two-column layout
                    self._add_two_column_slide(prs, slide_info, img_path, idx + 1, len(slides_data))
                else:
                    # High-impact content slide (respecting Rule of Six)
                    self._add_content_slide(prs, slide_info, idx + 1, len(slides_data))

            output_path = self._unique_path(os.path.join(self.output_dir, filename))
            prs.save(output_path)
            logger.info(f"Presentation saved to: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"Error creating presentation: {e}", exc_info=True)
            return None

    # ──────────────────────────────────────────────────────────────────
    # Slide builders
    # ──────────────────────────────────────────────────────────────────

    def _add_cover_slide(self, prs, slide_info, slide_num, total):
        """Full-bleed dark cover with large title and subtitle."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        W, H = prs.slide_width, prs.slide_height

        # Background
        self._set_bg(slide, self.COLOR_BG_ALT)

        # Left accent bar (thick)
        self._add_rect(slide, 0, 0, Inches(0.18), H, self.COLOR_ACCENT)

        # Decorative bottom bar
        self._add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), self.COLOR_ACCENT)

        # Title
        title_text = slide_info.get("title", "Presentation")
        tb = slide.shapes.add_textbox(Inches(0.55), Inches(1.8), Inches(12.0), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title_text
        run.font.name  = self.FONT_TITLE
        run.font.size  = Pt(52)
        run.font.bold  = True
        run.font.color.rgb = self.COLOR_TITLE

        # Subtitle (first bullet treated as subtitle)
        content = self._coerce_list(slide_info.get("content", []))
        subtitle = content[0] if content else ""
        if subtitle:
            tb2 = slide.shapes.add_textbox(Inches(0.55), Inches(4.0), Inches(10.0), Inches(0.9))
            tf2 = tb2.text_frame
            p2  = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.LEFT
            r2  = p2.add_run()
            r2.text = subtitle
            r2.font.name  = self.FONT_BODY
            r2.font.size  = Pt(22)
            r2.font.color.rgb = self.COLOR_BODY

        # Divider line between title and subtitle
        self._add_rect(slide, Inches(0.55), Inches(3.85), Inches(5.0), Inches(0.04), self.COLOR_ACCENT)

        # Slide counter
        self._add_slide_number(slide, prs, slide_num, total)
        self._add_speaker_notes(slide, slide_info)

    def _add_content_slide(self, prs, slide_info, slide_num, total):
        """Standard content slide with left accent bar and bullet list."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        W, H = prs.slide_width, prs.slide_height

        self._set_bg(slide, self.COLOR_BG)

        # Top accent bar
        self._add_rect(slide, 0, 0, W, Inches(0.08), self.COLOR_ACCENT)

        # Left accent bar
        self._add_rect(slide, 0, Inches(0.08), Inches(0.08), H - Inches(0.16), self.COLOR_ACCENT)

        # Title — centered across full width
        title_text = slide_info.get("title", "")
        tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.85))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title_text
        run.font.name  = self.FONT_TITLE
        run.font.size  = Pt(36)
        run.font.bold  = True
        run.font.color.rgb = self.COLOR_TITLE

        # Thin divider — centered under title, width scales with title length
        char_width_inches = 0.22
        divider_w = min(max(len(title_text) * char_width_inches, 1.5), 10.0)
        divider_left = (prs.slide_width.inches - divider_w) / 2
        self._add_rect(slide, Inches(divider_left), Inches(1.05), Inches(divider_w), Inches(0.02), self.COLOR_ACCENT)

        # Content bullets
        content = self._coerce_list(slide_info.get("content", []))
        if content:
            tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.0), Inches(6.0))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            for i, point in enumerate(content):
                para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                para.space_before = Pt(6)
                para.space_after  = Pt(4)
                # Bullet dot
                run_dot = para.add_run()
                run_dot.text = "◆  "
                run_dot.font.name  = self.FONT_TITLE
                run_dot.font.size  = Pt(10)
                run_dot.font.color.rgb = self.COLOR_ACCENT
                # Bullet text
                run_txt = para.add_run()
                run_txt.text = point
                run_txt.font.name  = self.FONT_BODY
                run_txt.font.size  = Pt(18)
                run_txt.font.color.rgb = self.COLOR_BODY

        self._add_slide_number(slide, prs, slide_num, total)
        self._add_speaker_notes(slide, slide_info)

    def _add_two_column_slide(self, prs, slide_info, img_path, slide_num, total):
        """Two-column slide: text left, image right.
        This method is only called when img_path is a valid, existing file."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W, H = prs.slide_width, prs.slide_height

        self._set_bg(slide, self.COLOR_BG)

        # Top accent bar
        self._add_rect(slide, 0, 0, W, Inches(0.08), self.COLOR_ACCENT)

        # Card background for right column
        self._add_rect(slide, Inches(7.2), Inches(0.08), Inches(6.133), H - Inches(0.08), self.COLOR_CARD_BG)

        # Title
        title_text = slide_info.get("title", "")
        tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(6.6), Inches(1.0))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.name  = self.FONT_TITLE
        run.font.size  = Pt(34)
        run.font.bold  = True
        run.font.color.rgb = self.COLOR_TITLE

        char_width_inches = 0.21
        divider_w = min(max(len(title_text) * char_width_inches, 1.5), 6.0)
        self._add_rect(slide, Inches(0.3), Inches(1.25), Inches(divider_w), Inches(0.02), self.COLOR_ACCENT)

        # Bullets (left column)
        content = self._coerce_list(slide_info.get("content", []))
        if content:
            tb2 = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(6.7), Inches(5.8))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            for i, point in enumerate(content):
                para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                para.space_before = Pt(7)
                para.space_after  = Pt(5)
                run_dot = para.add_run()
                run_dot.text = "◆  "
                run_dot.font.size  = Pt(9)
                run_dot.font.color.rgb = self.COLOR_ACCENT
                run_txt = para.add_run()
                run_txt.text = point
                run_txt.font.name  = self.FONT_BODY
                run_txt.font.size  = Pt(16)
                run_txt.font.color.rgb = self.COLOR_BODY

        # Insert user-provided image on the right column
        try:
            slide.shapes.add_picture(
                img_path,
                left=Inches(7.45), top=Inches(0.9),
                width=Inches(5.5),
            )
        except Exception as e:
            logger.error(f"Failed to insert image on slide {slide_num}: {e}")
            # Gracefully fall back — right column stays empty rather than crashing

        self._add_slide_number(slide, prs, slide_num, total)
        self._add_speaker_notes(slide, slide_info)

    def _add_closing_slide(self, prs, slide_info, slide_num, total):
        """Closing slide with gold accent — visually distinct from content slides."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W, H = prs.slide_width, prs.slide_height

        self._set_bg(slide, self.COLOR_BG_ALT)

        # Gold bottom bar
        self._add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), self.COLOR_ACCENT2)

        # Gold left bar
        self._add_rect(slide, 0, 0, Inches(0.18), H, self.COLOR_ACCENT2)

        # Large decorative circle (background element)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(8.5), Inches(-1.0), Inches(6.5), Inches(6.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(0x25, 0x32, 0x75)
        circle.line.fill.background()

        # Title
        title_text = slide_info.get("title", "Thank You")
        tb = slide.shapes.add_textbox(Inches(0.55), Inches(1.5), Inches(11.5), Inches(2.0))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title_text
        run.font.name  = self.FONT_TITLE
        run.font.size  = Pt(54)
        run.font.bold  = True
        run.font.color.rgb = self.COLOR_ACCENT2

        # Divider
        self._add_rect(slide, Inches(0.55), Inches(3.55), Inches(6.0), Inches(0.04), self.COLOR_ACCENT2)

        # Closing points
        content = self._coerce_list(slide_info.get("content", []))
        if content:
            tb2 = slide.shapes.add_textbox(Inches(0.55), Inches(3.8), Inches(11.5), Inches(3.2))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            for i, point in enumerate(content):
                para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                para.space_before = Pt(6)
                run_txt = para.add_run()
                run_txt.text = point
                run_txt.font.name  = self.FONT_BODY
                run_txt.font.size  = Pt(18)
                run_txt.font.color.rgb = self.COLOR_BODY

        self._add_slide_number(slide, prs, slide_num, total)
        self._add_speaker_notes(slide, slide_info)

    # ──────────────────────────────────────────────────────────────────
    # Helpers
    # ──────────────────────────────────────────────────────────────────

    def _set_bg(self, slide, color: RGBColor):
        bg   = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_rect(self, slide, left, top, width, height, color: RGBColor):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _add_slide_number(self, slide, prs, slide_num: int, total: int):
        W = prs.slide_width
        H = prs.slide_height
        tb = slide.shapes.add_textbox(W - Inches(1.3), H - Inches(0.45), Inches(1.1), Inches(0.35))
        tf = tb.text_frame
        p  = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{slide_num} / {total}"
        run.font.name  = self.FONT_BODY
        run.font.size  = Pt(10)
        run.font.color.rgb = self.COLOR_MUTED

    def _add_speaker_notes(self, slide, slide_info: dict):
        """Adds speaker notes and Creative Director visual instructions."""
        notes_text = slide_info.get("notes", "")
        
        # Append Visual Design Block if provided by the Architect
        visual = slide_info.get("visual")
        if visual and isinstance(visual, dict):
            visual_block = (
                f"\n\n🎨 VISUAL DESIGN (Creative Director)\n"
                f"--------------------------------------\n"
                f"IMAGE CONCEPT: {visual.get('concept', 'N/A')}\n\n"
                f"AI IMAGE PROMPT: {visual.get('prompt', 'N/A')}\n"
                f"STYLE: Clean 3D Isometric\n"
            )
            notes_text += visual_block

        if notes_text.strip():
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text.strip()

    @staticmethod
    def _coerce_list(content) -> list:
        if isinstance(content, str):
            return [content] if content else []
        return list(content) if content else []

    @staticmethod
    def _unique_path(path: str) -> str:
        if not os.path.exists(path):
            return path
        base, ext = os.path.splitext(path)
        counter = 1
        while os.path.exists(f"{base}_{counter}{ext}"):
            counter += 1
        return f"{base}_{counter}{ext}"